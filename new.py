
import os
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain.prompts import PromptTemplate
from langchain.chains.question_answering import load_qa_chain
from sentence_transformers import SentenceTransformer
from langchain.docstore.document import Document
from langchain_google_genai import ChatGoogleGenerativeAI
import google.generativeai as genai
from dotenv import load_dotenv
import time
from functools import lru_cache
import logging
from crewai.tools import tool

# Additional imports for multi-format extraction
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
if not GOOGLE_API_KEY:
    raise ValueError("GOOGLE_API_KEY not found in environment variables")

# Configure Gemini API
genai.configure(api_key=GOOGLE_API_KEY)

# Initialize the all-MiniLM-L6-v2 model for embeddings (free local model)
embedding_model = SentenceTransformer('all-mpnet-base-v2')

# Custom embedding function for LangChain
class MiniLMEmbeddings:
    def __init__(self, model):
        self.model = model

    def embed_documents(self, texts):
        """Generate embeddings for a list of texts."""
        logger.info("Generating embeddings for documents")
        return self.model.encode(texts, show_progress_bar=True).tolist()

    def embed_query(self, text):
        """Generate embedding for a single query."""
        logger.info("Generating embedding for query")
        return self.model.encode([text])[0].tolist()

# Cache query embeddings to reduce computation
@lru_cache(maxsize=100)
def get_cached_embedding(query):
    """Generate and cache query embedding using all-MiniLM-L6-v2."""
    embeddings = MiniLMEmbeddings(embedding_model)
    return embeddings.embed_query(query)

def get_conversational_chain():
    """Create QA chain with Google's Gemini model."""
    prompt_template = """
    Answer the question as detailed as possible from the provided context, make sure to provide all the details of rfp templates and sample rfp docs given to you, if the answer is not available in the provided context just say, "answer is not available in the context", don't provide the wrong answer.\n \n
    Context:\n{context}\n
    Question: \n{question}\n
    Answer:
    """
    try:
        logger.info("Creating conversational chain with Gemini")
        # Use Gemini 1.5 Pro (or 'gemini-1.5-flash' for faster/cheaper option)
        model = ChatGoogleGenerativeAI(model="gemini-2.5-flash", temperature=0.1)
        prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
        chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
        return chain
    except Exception as e:
        logger.error(f"Error creating conversational chain: {str(e)}")
        raise

# === Extraction helpers ===

def split_text_to_chunks(text: str):
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    return text_splitter.split_text(text)

def process_pdf(file_path):
    """Extract text from a single PDF and return text chunks."""
    try:
        logger.info(f"Processing PDF: {file_path}")
        reader = PdfReader(file_path)
        text = ''
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + '\n'
        if not text.strip():
            raise Exception(f"No text extracted from {file_path}")
        chunks = split_text_to_chunks(text)
        if not chunks:
            raise Exception(f"No text chunks created from {file_path}")
        logger.info(f"Extracted {len(chunks)} chunks from {file_path}")
        return chunks
    except Exception as e:
        logger.error(f"Error processing {file_path}: {str(e)}")
        return []

def process_docx(file_path: str):
    """Extract text from DOCX and return text chunks."""
    try:
        logger.info(f"Processing DOCX: {file_path}")
        doc = DocxDocument(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
        text = "\n".join(paragraphs)
        if not text.strip():
            raise Exception(f"No text extracted from {file_path}")
        chunks = split_text_to_chunks(text)
        logger.info(f"Extracted {len(chunks)} chunks from {file_path}")
        return chunks
    except Exception as e:
        logger.error(f"Error processing DOCX {file_path}: {str(e)}")
        return []

def process_xlsx(file_path: str):
    """Extract text from XLSX by reading cell values and returning chunks."""
    try:
        logger.info(f"Processing XLSX: {file_path}")
        # Use pandas to read all sheets
        xls = pd.ExcelFile(file_path)
        parts = []
        for sheet_name in xls.sheet_names:
            try:
                df = xls.parse(sheet_name=sheet_name, dtype=str)
                df = df.fillna("")
                # Create a textual representation: include sheet name
                parts.append(f"Sheet: {sheet_name}")
                parts.append(df.to_csv(index=False))
            except Exception as se:
                logger.warning(f"Failed parsing sheet {sheet_name} in {file_path}: {se}")
        text = "\n\n".join(parts)
        if not text.strip():
            raise Exception(f"No text extracted from {file_path}")
        chunks = split_text_to_chunks(text)
        logger.info(f"Extracted {len(chunks)} chunks from {file_path}")
        return chunks
    except Exception as e:
        logger.error(f"Error processing XLSX {file_path}: {str(e)}")
        return []

def process_pptx(file_path: str):
    """Extract text from PPTX slides and return chunks."""
    try:
        logger.info(f"Processing PPTX: {file_path}")
        prs = Presentation(file_path)
        parts = []
        for idx, slide in enumerate(prs.slides, start=1):
            slide_text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text_parts.append(shape.text)
            if slide_text_parts:
                parts.append(f"Slide {idx}:\n" + "\n".join(slide_text_parts))
        text = "\n\n".join(parts)
        if not text.strip():
            raise Exception(f"No text extracted from {file_path}")
        chunks = split_text_to_chunks(text)
        logger.info(f"Extracted {len(chunks)} chunks from {file_path}")
        return chunks
    except Exception as e:
        logger.error(f"Error processing PPTX {file_path}: {str(e)}")
        return []

def process_file_by_extension(file_path: str):
    ext = os.path.splitext(file_path.lower())[1]
    if ext == '.pdf':
        return process_pdf(file_path)
    if ext == '.docx':
        return process_docx(file_path)
    if ext == '.xlsx':
        return process_xlsx(file_path)
    if ext == '.pptx':
        return process_pptx(file_path)
    logger.warning(f"Unsupported file type for {file_path}; skipping.")
    return []

def load_or_create_vector_db(folder_path="uploads", persist_directory="db"):
    """Load existing Chroma store or create a new one from supported files in the folder."""
    embeddings = MiniLMEmbeddings(embedding_model)
    try:
        if os.path.exists(os.path.join(persist_directory, "chroma.sqlite3")):
            logger.info(f"Loading existing Chroma store from {persist_directory}")
            return Chroma(
                persist_directory=persist_directory,
                embedding_function=embeddings
            )

        if not os.path.exists(folder_path):
            raise Exception(f"Folder {folder_path} does not exist")

        supported_exts = {'.pdf', '.docx', '.xlsx', '.pptx'}
        files = [f for f in os.listdir(folder_path) if os.path.splitext(f.lower())[1] in supported_exts]
        if not files:
            raise Exception(f"No supported files found in {folder_path}")

        logger.info(f"Processing {len(files)} files from {folder_path}")
        all_chunks = []
        for fname in files:
            file_path = os.path.join(folder_path, fname)
            logger.info(f"Extracting from: {file_path}")
            chunks = process_file_by_extension(file_path)
            if chunks:
                all_chunks.extend(chunks)

        if not all_chunks:
            raise Exception("No valid text chunks extracted from any document")

        logger.info(f"Total chunks to embed: {len(all_chunks)}")
        logger.info("Generating embeddings... This may take a while")

        documents = [Document(page_content=chunk) for chunk in all_chunks]

        db = Chroma.from_documents(
            documents=documents,
            embedding=embeddings,
            persist_directory=persist_directory
        )

        logger.info(f"Created and saved new Chroma store to {persist_directory}")
        return db

    except Exception as e:
        logger.error(f"Error loading or creating vector store: {str(e)}")
        raise

@tool("ChromaDB Query Tool")
def query_vector_db1(query: str) -> str:
    """Searches through the rfp instruction and company data ChromaDB vector database and provides a response to the given query using AI for QA based on LangChain."""
    try:
        logger.info(f"Querying vector store with: {query}")
        db = load_or_create_vector_db()
        docs = db.similarity_search(query, k=20)
        if not docs:
            logger.info("No relevant documents found")
            return "No relevant information found in the vector store."

        chain = get_conversational_chain()
        response = chain(
            {"input_documents": docs, "question": query},
            return_only_outputs=True
        )
        logger.info("Query processed successfully")
        return response["output_text"]
    except Exception as e:
        logger.error(f"Error querying vector store: {str(e)}")
        return f"Error: {str(e)}"

# Optional: quick smoke test if this file is executed directly
if __name__ == "__main__":
    db = load_or_create_vector_db()
    docs = db.similarity_search("what's the data is about in documents?", k=20)
    if not docs:
        logger.info("No relevant documents found")
        print("No relevant information found in the vector store.")
    else:
        chain = get_conversational_chain()
        response = chain(
            {"input_documents": docs, "question": "what's the data is about in documents?"},
            return_only_outputs=True
        )
        print(response["output_text"])

